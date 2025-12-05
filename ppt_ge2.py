from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # 创建演示文稿对象
    prs = Presentation()

    # --- 幻灯片 1: 封面 ---
    slide_layout = prs.slide_layouts[0] # 标题幻灯片布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "自然灾害指数保险\n端到端解决方案"
    subtitle.text = "基于 Google Map Platform 与 AI 赋能的普惠型保险服务平台\n\n事前预警防范 • 事后快速赔付"

    # --- 幻灯片 2: 核心目标 (Goal) ---
    slide_layout = prs.slide_layouts[1] # 标题+内容布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "核心目标 (Goal)"
    
    # 获取内容文本框
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear() # 清除默认格式

    # 添加要点
    points = [
        ("端到端防护", "提供降低自然灾害影响的完整方案，实现“事前预警防范 + 事后快速赔付”闭环。"),
        ("全覆盖普惠服务", "服务各阶层用户，特别覆盖低文化水平及使用功能机（非智能机）的用户群体。"),
        ("直观信息传递", "降低沟通门槛，提供更有说服力的公开数据和直观的风险评估。")
    ]

    for bold_text, normal_text in points:
        p = tf.add_paragraph()
        p.text = f"{bold_text}：{normal_text}"
        p.level = 0
        p.space_after = Pt(14)

    # --- 幻灯片 3: 建设任务 (Objective) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "建设任务 (Objective)"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()

    # 主点 1
    p = tf.add_paragraph()
    p.text = "统一 Index Insurance 门户"
    p.level = 0
    p.font.bold = True

    sub_points = [
        "集合多种创新的 Index 保险产品",
        "集成历史、当前和预测的天气数据，叠加产品逻辑进行可视化",
        "情景式 AI 咨询窗口，提供实时交互"
    ]
    for sp in sub_points:
        p = tf.add_paragraph()
        p.text = sp
        p.level = 1

    # 主点 2 & 3
    other_points = [
        ("新分销形态", "AI 电销：利用人工智能进行主动触达与产品推广"),
        ("新售后形态", "AI 电话服务：覆盖预警通知、保单管理及理赔全流程")
    ]

    for bold_text, normal_text in other_points:
        p = tf.add_paragraph()
        p.text = f"{bold_text} —— {normal_text}"
        p.level = 0
        p.space_before = Pt(14)


    # --- 幻灯片 4: 核心策略 (Core Approach) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "核心策略 (Core Approach)"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()

    strategies = [
        ("深度集成 Google Map Platform", "全面涵盖 Weather API, MCP (Mission Control Platform) 和 AI Kit，构建坚实技术底座。"),
        ("数据驱动的预警机制", "基于精准的天气预测数据，建立自动化的预警服务与响应机制，变被动为主动。"),
        ("Multi-Agent 多级合作模式", "采用多智能体协作架构，各司其职，高效处理复杂的保险业务与灾害应对流程。")
    ]

    for title_text, desc_text in strategies:
        p = tf.add_paragraph()
        p.text = title_text
        p.font.bold = True
        p.level = 0
        
        p2 = tf.add_paragraph()
        p2.text = desc_text
        p2.level = 1
        p2.space_after = Pt(14)

    # --- 幻灯片 5: 总结 ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "总结：构建更有韧性的未来"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()

    keywords = "关键词：科技普惠  |  精准风控  |  智能协作"
    p = tf.add_paragraph()
    p.text = keywords
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(24)

    summary_text = "通过 Google Maps 强大的地理数据能力与先进的 AI Agent 架构，我们致力于为每一位用户提供公平、透明、及时的自然灾害保险保障。"
    p = tf.add_paragraph()
    p.text = summary_text
    p.level = 0

    # 保存文件
    output_filename = 'Natural_Disaster_Insurance_Presentation.pptx'
    prs.save(output_filename)
    print(f"成功生成演示文稿: {output_filename}")

if __name__ == "__main__":
    create_presentation()