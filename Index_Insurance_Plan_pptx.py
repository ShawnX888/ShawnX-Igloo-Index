from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Initialize Presentation
    prs = Presentation()

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI-Driven Index Insurance Platform"
    subtitle.text = "End-to-End Disaster Resilience & Inclusive Protection"

    # --- Slide 2: Goal ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Strategic Goals"
    text_frame = content.text_frame
    text_frame.text = "End-to-End Solution"
    
    # Add bullets
    p = text_frame.add_paragraph()
    p.text = "Mitigate disaster impact: Pre-warning prevention + Post-disaster compensation."
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "Inclusive Access"
    p = text_frame.add_paragraph()
    p.text = "Serve users of all levels, including low-literacy & feature phone users."
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Lower Communication Threshold"
    p = text_frame.add_paragraph()
    p.text = "More intuitive information delivery and visualized data."
    p.level = 1

    p = text_frame.add_paragraph()
    p.text = "Trustworthy Risk Assessment"
    p = text_frame.add_paragraph()
    p.text = "Based on convincing public data and transparent evaluation."
    p.level = 1

    # --- Slide 3: Objective ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Objectives"
    text_frame = content.text_frame
    text_frame.text = "Build a Unified Index Insurance Portal"
    
    # Sub-bullets for Portal
    p = text_frame.add_paragraph()
    p.text = "Aggregate innovative index insurance products."
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "Integrate historical, current, and forecasted weather data."
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "Visualize logic: Overlay Index products onto weather data."
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "Scenario-based AI consultation window."
    p.level = 1

    # Other Objectives
    p = text_frame.add_paragraph()
    p.text = "Introduce New Distribution Forms"
    p = text_frame.add_paragraph()
    p.text = "AI-driven E-commerce and Telemarketing."
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "Introduce New After-Sales Forms"
    p = text_frame.add_paragraph()
    p.text = "AI Telephony Service (Warnings, Policy Management & Claims)."
    p.level = 1

    # --- Slide 4: Core Approach ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Core Approach"
    text_frame = content.text_frame
    
    p = text_frame.add_paragraph()
    p.text = "Deep Integration with Google Maps Platform"
    p = text_frame.add_paragraph()
    p.text = "Leveraging Weather API, MCP, and AI Kit."
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "Proactive Warning Mechanism"
    p = text_frame.add_paragraph()
    p.text = "Establishing warning services based on precise weather forecasts."
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "Multi-Agent Collaboration"
    p = text_frame.add_paragraph()
    p.text = "Multi-level cooperation model where each agent performs specific duties."
    p.level = 1

    # Save the presentation
    file_name = "Index_Insurance_Plan.pptx"
    prs.save(file_name)
    print(f"Successfully generated {file_name}")

if __name__ == "__main__":
    create_presentation()